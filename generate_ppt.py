"""
AgriPredict AI — Full Project Presentation Generator
Generates a professional PPT + converts to PDF
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os

# ── Color Palette ──────────────────────────────────────────────
GREEN      = RGBColor(0x16, 0xA3, 0x4A)   # #16a34a
DARK_GREEN = RGBColor(0x14, 0x53, 0x2D)   # #14532d
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
DARK_BG    = RGBColor(0x0F, 0x17, 0x2A)   # #0f172a
SLATE      = RGBColor(0x1E, 0x29, 0x3B)   # #1e293b
ACCENT     = RGBColor(0x22, 0xC5, 0x5E)   # #22c55e
MUTED      = RGBColor(0x64, 0x74, 0x8B)   # #64748b

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helper: fill slide background ─────────────────────────────
def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# ── Helper: add a text box ─────────────────────────────────────
def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb

# ── Helper: add rectangle shape ───────────────────────────────
def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

# ── Helper: add bullet list ────────────────────────────────────
def add_bullets(slide, items, left, top, width, height,
                font_size=16, color=LIGHT_GRAY, indent=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        bullet = "  •  " if indent else "• "
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"

# ── Helper: slide header strip ────────────────────────────────
def add_header(slide, title, subtitle=None):
    # Green top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), GREEN)
    # Title
    add_text(slide, title,
             Inches(0.4), Inches(0.08), Inches(12), Inches(0.65),
             font_size=28, bold=True, color=WHITE, font_name="Calibri")
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.72), Inches(12), Inches(0.35),
                 font_size=14, bold=False, color=LIGHT_GRAY, font_name="Calibri")

# ── Helper: footer ────────────────────────────────────────────
def add_footer(slide, text="AgriPredict AI  |  MCA Final Year Project  |  2026"):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), DARK_GREEN)
    add_text(slide, text,
             Inches(0.3), Inches(7.17), Inches(13), Inches(0.28),
             font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  SLIDE CREATORS
# ══════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title / Cover"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(sl, DARK_BG)

    # Left decorative bar
    add_rect(sl, 0, 0, Inches(0.18), SLIDE_H, GREEN)

    # Green accent line (center)
    add_rect(sl, Inches(0.5), Inches(3.35), Inches(5.5), Inches(0.05), ACCENT)

    # Main title
    add_text(sl, "AgriPredict AI",
             Inches(0.5), Inches(1.2), Inches(12), Inches(1.0),
             font_size=52, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri")

    # Subtitle
    add_text(sl,
             "AI-Driven Fruits & Vegetables Price Forecast\n& Smart Market Recommendation System",
             Inches(0.5), Inches(2.3), Inches(11), Inches(1.0),
             font_size=22, bold=False, color=ACCENT,
             align=PP_ALIGN.LEFT, font_name="Calibri")

    # Details
    add_text(sl,
             "MCA Final Year Project  |  2026\nDepartment of Computer Applications",
             Inches(0.5), Inches(3.6), Inches(9), Inches(0.8),
             font_size=15, color=MUTED, font_name="Calibri")

    # Tag chips
    tags = ["Python Flask", "XGBoost ML", "SQLite", "Bootstrap 5", "Chart.js"]
    x = Inches(0.5)
    for tag in tags:
        add_rect(sl, x, Inches(4.7), Inches(1.7), Inches(0.38), SLATE)
        add_text(sl, tag, x + Inches(0.08), Inches(4.72),
                 Inches(1.55), Inches(0.34),
                 font_size=11, color=ACCENT, align=PP_ALIGN.CENTER)
        x += Inches(1.85)

    add_footer(sl)
    return sl


def slide_problem(prs):
    """Slide 2 — Problem Statement & Objectives"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "Problem Statement & Objectives",
               "Why AgriPredict AI was built")
    add_footer(sl)

    # Left column — Problem
    add_rect(sl, Inches(0.3), Inches(1.3), Inches(5.8), Inches(5.5), SLATE)
    add_text(sl, "⚠  The Problem",
             Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.45),
             font_size=16, bold=True, color=ACCENT)
    problems = [
        "Farmers earn 33–60% less due to unpredictable price volatility",
        "No accessible tool to forecast next-week prices at a district level",
        "Weather events (rainfall, temperature) cause sudden price swings",
        "Traders and intermediaries exploit information asymmetry",
        "Farmers sell at wrong time, wrong market — losing profit",
    ]
    add_bullets(sl, problems, Inches(0.5), Inches(1.95), Inches(5.5), Inches(4.5),
                font_size=14, color=LIGHT_GRAY)

    # Right column — Objectives
    add_rect(sl, Inches(6.4), Inches(1.3), Inches(6.5), Inches(5.5), SLATE)
    add_text(sl, "🎯  Project Objectives",
             Inches(6.6), Inches(1.4), Inches(6.2), Inches(0.45),
             font_size=16, bold=True, color=ACCENT)
    objectives = [
        "Forecast next-week modal prices for 19 agri commodities",
        "Use XGBoost ML with weather + location features",
        "Recommend the most profitable nearby market district",
        "Provide actionable Hold vs Sell advisory to farmers",
        "Deliver a multilingual (5 languages) web interface",
        "Achieve fast predictions with <2 second response time",
    ]
    add_bullets(sl, objectives, Inches(6.6), Inches(1.95), Inches(6.1), Inches(4.5),
                font_size=14, color=LIGHT_GRAY)
    return sl


def slide_architecture(prs):
    """Slide 3 — System Architecture"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "System Architecture", "End-to-end flow of the application")
    add_footer(sl)

    layers = [
        ("🌐  Frontend Layer",
         "Bootstrap 5 • Custom CSS Glassmorphism • Chart.js • i18n JS engine"),
        ("🔐  Authentication Layer",
         "Flask-Login • Session management • SQLite User DB • Secure password storage"),
        ("📊  Dashboard & Input Layer",
         "Commodity / State / District dropdowns • Dynamic district loading via JS"),
        ("🤖  ML Inference Layer",
         "XGBoost Regressor model • Label-encoded features • Joblib model loading"),
        ("🌦  Weather Context Layer",
         "Mock weather engine (temp 25–35°C, rainfall 10–50mm, humidity 40–70%)"),
        ("📈  Output & Advisory Layer",
         "Price forecast • Trend detection • Market comparison • Farmer advisory"),
    ]

    y = Inches(1.25)
    for i, (title, desc) in enumerate(layers):
        color = SLATE if i % 2 == 0 else RGBColor(0x1A, 0x31, 0x4A)
        add_rect(sl, Inches(0.3), y, Inches(12.7), Inches(0.82), color)
        # Arrow connector
        if i < len(layers) - 1:
            add_rect(sl, Inches(6.4), y + Inches(0.82), Inches(0.05), Inches(0.12), ACCENT)
        add_text(sl, title,
                 Inches(0.5), y + Inches(0.05), Inches(3.8), Inches(0.4),
                 font_size=14, bold=True, color=ACCENT)
        add_text(sl, desc,
                 Inches(4.4), y + Inches(0.05), Inches(8.4), Inches(0.7),
                 font_size=13, color=LIGHT_GRAY)
        y += Inches(0.95)
    return sl


def slide_ml(prs):
    """Slide 4 — Machine Learning Implementation"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "Machine Learning Implementation",
               "XGBoost model training, features, and evaluation")
    add_footer(sl)

    # Dataset card
    add_rect(sl, Inches(0.3), Inches(1.3), Inches(4.0), Inches(5.5), SLATE)
    add_text(sl, "📂  Dataset",
             Inches(0.5), Inches(1.4), Inches(3.7), Inches(0.4),
             font_size=15, bold=True, color=ACCENT)
    ds_items = [
        "50,000+ synthetic records",
        "19 commodities",
        "5 Indian states",
        "Multiple districts",
        "Features:",
        "  Month, Commodity",
        "  State, District",
        "  Temperature (°C)",
        "  Rainfall (mm)",
        "  Humidity (%)",
        "Target: Modal Price/kg",
    ]
    add_bullets(sl, ds_items, Inches(0.5), Inches(1.9), Inches(3.6), Inches(4.8),
                font_size=12.5, color=LIGHT_GRAY, indent=False)

    # Model card
    add_rect(sl, Inches(4.55), Inches(1.3), Inches(4.0), Inches(5.5), SLATE)
    add_text(sl, "🤖  Model Selection",
             Inches(4.75), Inches(1.4), Inches(3.7), Inches(0.4),
             font_size=15, bold=True, color=ACCENT)
    model_items = [
        "Models evaluated:",
        "  Linear Regression",
        "  Random Forest",
        "  ✅ XGBoost Regressor",
        "",
        "XGBoost selected because:",
        "  Handles non-linear data",
        "  Best R² score",
        "  Robust to outliers",
        "  Fast inference time",
        "  Handles missing vals",
    ]
    add_bullets(sl, model_items, Inches(4.75), Inches(1.9), Inches(3.6), Inches(4.8),
                font_size=12.5, color=LIGHT_GRAY, indent=False)

    # Pipeline card
    add_rect(sl, Inches(8.8), Inches(1.3), Inches(4.2), Inches(5.5), SLATE)
    add_text(sl, "⚙  ML Pipeline",
             Inches(9.0), Inches(1.4), Inches(3.9), Inches(0.4),
             font_size=15, bold=True, color=ACCENT)
    pipe_items = [
        "1. Data loading (CSV)",
        "2. Label Encoding",
        "   Commodity, State,",
        "   District columns",
        "3. Train/Test split",
        "   (80/20 ratio)",
        "4. XGBoost training",
        "5. Joblib model save",
        "   xgboost_model.pkl",
        "   label_encoders.pkl",
        "6. Flask inference",
    ]
    add_bullets(sl, pipe_items, Inches(9.0), Inches(1.9), Inches(3.8), Inches(4.8),
                font_size=12.5, color=LIGHT_GRAY, indent=False)
    return sl


def slide_tech(prs):
    """Slide 5 — Technology Stack"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "Technology Stack", "Full-stack tools and frameworks used")
    add_footer(sl)

    stacks = [
        ("Frontend",   "🌐", ["HTML5 & CSS3 (Glassmorphism)", "Bootstrap 5", "JavaScript (ES6+)", "Chart.js (Price Trend Graph)", "i18n Engine (5 languages)"]),
        ("Backend",    "⚙",  ["Python 3.x", "Flask Web Framework", "Flask-Login (Auth)", "Flask-SQLAlchemy (ORM)", "Jinja2 Templating"]),
        ("ML / Data",  "🤖", ["Pandas & NumPy", "Scikit-learn", "XGBoost Regressor", "Joblib (Model I/O)", "Custom Advisory Engine"]),
        ("Database",   "🗄",  ["SQLite (via SQLAlchemy)", "agricultural_db.sqlite3", "User table (id, name,", " email, password)", "Lightweight & portable"]),
    ]

    x = Inches(0.3)
    for (title, icon, items) in stacks:
        add_rect(sl, x, Inches(1.3), Inches(3.0), Inches(5.5), SLATE)
        add_text(sl, icon + "  " + title,
                 x + Inches(0.15), Inches(1.4), Inches(2.7), Inches(0.45),
                 font_size=15, bold=True, color=ACCENT)
        add_bullets(sl, items, x + Inches(0.1), Inches(1.95),
                    Inches(2.8), Inches(4.5), font_size=13, color=LIGHT_GRAY, indent=False)
        x += Inches(3.25)
    return sl


def slide_features(prs):
    """Slide 6 — Key Features"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "Key Features & Modules",
               "What AgriPredict AI offers to farmers and traders")
    add_footer(sl)

    features = [
        ("🔐 Secure Auth",
         "User registration & login with Flask-Login session management. Email-based identity."),
        ("🌾 19 Commodities",
         "Tomato, Onion, Potato, Banana, Apple, Mango, Cabbage, Carrot, Brinjal, Grapes, Orange, Papaya, Pomegranate, Watermelon, Spinach, Cauli., Garlic, Ginger, Green Chilli."),
        ("📍 Location-Based",
         "5 states × 5 districts each. District-level precision for AP, Telangana, Karnataka, Tamil Nadu, Maharashtra."),
        ("🌦 Weather-Aware",
         "Temperature, rainfall, humidity integrated into price forecasting for realistic predictions."),
        ("📊 Price Chart",
         "Interactive Chart.js graph shows historical trend (past 8 weeks) vs forecasted price for next week."),
        ("🏪 Market Finder",
         "Compares selected district vs 3 nearby districts, recommends highest-price market automatically."),
        ("💡 Farmer Advisory",
         "Dynamic Hold vs Sell recommendations with weather-context reasoning. Multilingual (EN/TE/HI/TA/KN)."),
        ("🌐 Multilingual UI",
         "Full i18n support for English, Telugu, Hindi, Tamil, Kannada via JSON language packs."),
    ]

    x, y = Inches(0.3), Inches(1.3)
    for i, (title, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        lx = Inches(0.3) + col * Inches(6.5)
        ly = Inches(1.3) + row * Inches(1.42)
        add_rect(sl, lx, ly, Inches(6.2), Inches(1.3), SLATE)
        add_text(sl, title, lx + Inches(0.15), ly + Inches(0.07),
                 Inches(5.9), Inches(0.38),
                 font_size=13, bold=True, color=ACCENT)
        add_text(sl, desc, lx + Inches(0.15), ly + Inches(0.48),
                 Inches(5.9), Inches(0.78),
                 font_size=11.5, color=LIGHT_GRAY)
    return sl


def slide_modules(prs):
    """Slide 7 — Module Workflow"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "Application Workflow",
               "Step-by-step user journey through the system")
    add_footer(sl)

    steps = [
        ("1", "Register / Login",     "User creates account or signs in securely"),
        ("2", "Select Parameters",    "Choose commodity, state & district from dashboard"),
        ("3", "Weather Data",         "System fetches / generates weather context data"),
        ("4", "ML Prediction",        "XGBoost model returns next-week forecasted price"),
        ("5", "Market Comparison",    "Model scans 3 nearby districts for best price"),
        ("6", "Advisory Generated",   "System outputs Hold/Sell advice with reasoning"),
        ("7", "Results Displayed",    "Chart, price, market, advisory shown on results page"),
    ]

    y = Inches(1.3)
    for (num, title, desc) in steps:
        # number circle
        add_rect(sl, Inches(0.3), y, Inches(0.5), Inches(0.7), GREEN)
        add_text(sl, num, Inches(0.3), y + Inches(0.05), Inches(0.5), Inches(0.58),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # card
        add_rect(sl, Inches(0.95), y, Inches(12.05), Inches(0.7), SLATE)
        add_text(sl, title, Inches(1.1), y + Inches(0.05), Inches(3.5), Inches(0.35),
                 font_size=14, bold=True, color=ACCENT)
        add_text(sl, desc, Inches(4.8), y + Inches(0.05), Inches(7.9), Inches(0.6),
                 font_size=13, color=LIGHT_GRAY)
        # connector arrow
        if num != "7":
            add_rect(sl, Inches(0.49), y + Inches(0.7), Inches(0.12), Inches(0.12), MUTED)
        y += Inches(0.83)
    return sl


def slide_results(prs):
    """Slide 8 — Output & Results"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "Output & Results", "What the user sees after prediction")
    add_footer(sl)

    sections = [
        ("📈  Price Forecast Card",
         ["Current estimated modal price (₹/kg)", "AI forecasted price for next week",
          "% change with trend indicator (↑ or ↓)", "Temperature, Rainfall, Humidity shown"]),
        ("📊  Interactive Price Chart",
         ["Chart.js line graph", "8 weeks of historical prices plotted",
          "Forecasted price highlighted in green", "Smooth animation on page load"]),
        ("🏪  Market Recommendation",
         ["Table of 4 districts compared", "Prices sorted highest to lowest",
          "Best market highlighted in green", "Recommended district clearly labeled"]),
        ("💡  Farmer Advisory",
         ["Dynamic advisory based on ML output", "Hold vs Sell decision in plain language",
          "Weather-context reasoning included", "Fully translated in 5 Indian languages"]),
    ]

    x = Inches(0.3)
    for i, (title, bullets) in enumerate(sections):
        col = i % 2
        row = i // 2
        lx = Inches(0.3) + col * Inches(6.5)
        ly = Inches(1.3) + row * Inches(2.9)
        add_rect(sl, lx, ly, Inches(6.2), Inches(2.7), SLATE)
        add_rect(sl, lx, ly, Inches(6.2), Inches(0.45), GREEN)
        add_text(sl, title, lx + Inches(0.15), ly + Inches(0.04),
                 Inches(5.9), Inches(0.38),
                 font_size=13, bold=True, color=WHITE)
        add_bullets(sl, bullets, lx + Inches(0.1), ly + Inches(0.55),
                    Inches(5.9), Inches(2.0), font_size=12.5, color=LIGHT_GRAY, indent=True)
    return sl


def slide_database(prs):
    """Slide 9 — Database & File Structure"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "Database & Project Structure",
               "How data is stored and the codebase is organized")
    add_footer(sl)

    # DB schema
    add_rect(sl, Inches(0.3), Inches(1.3), Inches(5.5), Inches(5.5), SLATE)
    add_text(sl, "🗄  SQLite Database Schema",
             Inches(0.5), Inches(1.4), Inches(5.2), Inches(0.4),
             font_size=15, bold=True, color=ACCENT)
    db_items = [
        "users table:",
        "  id          INTEGER PRIMARY KEY",
        "  name        VARCHAR(100)",
        "  email       VARCHAR(100) UNIQUE",
        "  password    VARCHAR(100)",
        "",
        "File: instance/agricultural_db.sqlite3",
        "ORM: Flask-SQLAlchemy",
        "Created: Auto on first run (db.create_all)",
    ]
    add_bullets(sl, db_items, Inches(0.5), Inches(1.9), Inches(5.1), Inches(4.5),
                font_size=13, color=LIGHT_GRAY, indent=False)

    # File structure
    add_rect(sl, Inches(6.1), Inches(1.3), Inches(6.9), Inches(5.5), SLATE)
    add_text(sl, "📁  Project File Structure",
             Inches(6.3), Inches(1.4), Inches(6.6), Inches(0.4),
             font_size=15, bold=True, color=ACCENT)
    files = [
        "mca_project/",
        "  app.py              ← Main Flask app",
        "  generate_ppt.py     ← PPT generator",
        "  requirements.txt    ← Dependencies",
        "  agricultural_data.csv",
        "  models/",
        "    xgboost_model.pkl",
        "    label_encoders.pkl",
        "  templates/  (5 HTML files)",
        "    base, dashboard, login,",
        "    register, results",
        "  static/",
        "    css/style.css",
        "    js/i18n.js",
        "    lang/  en, te, hi, ta, kn .json",
        "  instance/  (SQLite DB)",
    ]
    add_bullets(sl, files, Inches(6.3), Inches(1.9), Inches(6.5), Inches(4.5),
                font_size=12, color=LIGHT_GRAY, indent=False)
    return sl


def slide_multilingual(prs):
    """Slide 10 — Multilingual Support"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "Multilingual Support (i18n)",
               "Complete 5-language translation system for Indian farmers")
    add_footer(sl)

    langs = [
        ("🇬🇧", "English",  "en.json", "Default language"),
        ("🇮🇳", "Telugu",   "te.json", "తెలుగు — Andhra & Telangana"),
        ("🇮🇳", "Hindi",    "hi.json", "हिंदी — North India"),
        ("🇮🇳", "Tamil",    "ta.json", "தமிழ் — Tamil Nadu"),
        ("🇮🇳", "Kannada",  "kn.json", "ಕನ್ನಡ — Karnataka"),
    ]

    x = Inches(0.4)
    for flag, lang, file, note in langs:
        add_rect(sl, x, Inches(1.4), Inches(2.3), Inches(1.1), SLATE)
        add_text(sl, flag + "  " + lang,
                 x + Inches(0.1), Inches(1.45), Inches(2.1), Inches(0.45),
                 font_size=16, bold=True, color=WHITE)
        add_text(sl, file, x + Inches(0.1), Inches(1.88),
                 Inches(2.1), Inches(0.3),
                 font_size=11, color=ACCENT)
        x += Inches(2.5)

    add_text(sl, "Note:  " + langs[1][3],
             Inches(0.4), Inches(2.65), Inches(12), Inches(0.3),
             font_size=12, color=MUTED)

    # How it works
    add_rect(sl, Inches(0.3), Inches(3.1), Inches(12.7), Inches(3.7), SLATE)
    add_text(sl, "⚙  How the i18n Engine Works",
             Inches(0.5), Inches(3.2), Inches(12), Inches(0.4),
             font_size=15, bold=True, color=ACCENT)
    how_items = [
        "Language JSON files stored in static/lang/ — separate file per language",
        "i18n.js engine loads correct JSON on page load based on user's selected language",
        "HTML elements use data-i18n attribute — JS replaces text content automatically",
        "Language preference saved to localStorage — persists across page navigations",
        "Navbar language switcher (🌐 dropdown) lets users switch in one click",
        "All UI labels, form placeholders, buttons, advisories — fully translated",
        "Advisory text is dynamically constructed in JS using translated phrase templates",
    ]
    add_bullets(sl, how_items, Inches(0.5), Inches(3.7), Inches(12.3), Inches(3.0),
                font_size=13, color=LIGHT_GRAY)
    return sl


def slide_advisory(prs):
    """Slide 11 — Farmer Advisory Engine"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "Farmer Advisory Engine",
               "Data-driven Hold vs Sell recommendations")
    add_footer(sl)

    add_rect(sl, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.5), SLATE)
    add_text(sl, "How the Advisory is Generated:",
             Inches(0.5), Inches(1.4), Inches(12), Inches(0.4),
             font_size=15, bold=True, color=ACCENT)

    logic = [
        "1.  Trend Detection        →   If predicted_price > current_price → 'Increasing'; else 'Decreasing'",
        "2.  Severity Classification →  ≥15% change = 'significantly' | ≥7% = 'moderately' | else 'slightly'",
        "3.  Commodity Classification →",
        "       Rain-sensitive: Tomato, Onion, Spinach, Green Chilli, Cabbage, Cauliflower",
        "       Heat-sensitive: Spinach, Cabbage, Cauliflower, Grapes",
        "       Heat-loving:    Watermelon, Mango, Banana, Papaya",
        "       Storable:       Potato, Onion, Garlic, Ginger (can hold stock safely)",
        "4.  Weather Context Clause →   Adds reasons like 'heavy rainfall disrupting transportation'",
        "5.  Action Advice →",
        "       Increasing + Perishable:  'Hold and sell next week for higher returns'",
        "       Increasing + Storable:    'Store safely and wait for price peak before selling'",
        "       Decreasing + Heavy:       'Sell immediately to avoid further losses'",
        "       Decreasing + Normal:      'Sell in batches rather than waiting'",
        "6.  Output →  Full English advisory + structured data for multilingual JS templates",
    ]
    add_bullets(sl, logic, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.7),
                font_size=12, color=LIGHT_GRAY, indent=False)
    return sl


def slide_conclusion(prs):
    """Slide 12 — Conclusion & Future Work"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_header(sl, "Conclusion & Future Work",
               "Summary and planned enhancements")
    add_footer(sl)

    # Conclusion
    add_rect(sl, Inches(0.3), Inches(1.3), Inches(12.7), Inches(2.4), SLATE)
    add_text(sl, "✅  Conclusion",
             Inches(0.5), Inches(1.4), Inches(12), Inches(0.4),
             font_size=15, bold=True, color=ACCENT)
    conclusions = [
        "Successfully built a full-stack AI web application for agricultural price forecasting",
        "XGBoost model trained on 50,000+ records with weather, location, and commodity features",
        "Covers 19 commodities across 5 Indian states with district-level precision",
        "Multilingual support (5 languages) makes it accessible to farmers across India",
        "Smart market recommendation & advisory engine helps maximize farmer profit",
    ]
    add_bullets(sl, conclusions, Inches(0.5), Inches(1.88), Inches(12.3), Inches(1.7),
                font_size=13, color=LIGHT_GRAY)

    # Future
    add_rect(sl, Inches(0.3), Inches(3.9), Inches(12.7), Inches(2.9), SLATE)
    add_text(sl, "🚀  Future Enhancements",
             Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
             font_size=15, bold=True, color=ACCENT)
    future = [
        "Live e-NAM / Agmarknet government API integration for real market prices",
        "Open-Meteo live weather API replacing mock weather data",
        "SMS/WhatsApp alert system for sudden price changes (Twilio integration)",
        "Mobile-first Progressive Web App (PWA) for offline rural access",
        "Expand to 100+ commodities and all Indian states / UTs",
        "Crop planting season advisor — recommend what to grow based on predicted demand",
    ]
    add_bullets(sl, future, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.1),
                font_size=13, color=LIGHT_GRAY)
    return sl


def slide_thankyou(prs):
    """Slide 13 — Thank You"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, DARK_BG)
    add_rect(sl, 0, 0, Inches(0.18), SLIDE_H, GREEN)

    add_text(sl, "Thank You",
             Inches(0.5), Inches(1.8), Inches(12), Inches(1.2),
             font_size=60, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font_name="Calibri")

    add_rect(sl, Inches(3.5), Inches(3.2), Inches(6.3), Inches(0.05), ACCENT)

    add_text(sl, "AgriPredict AI  —  MCA Final Year Project  |  2026",
             Inches(0.5), Inches(3.4), Inches(12), Inches(0.5),
             font_size=16, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(sl,
             "Built with  Python Flask  •  XGBoost  •  Bootstrap 5  •  Chart.js  •  SQLite",
             Inches(0.5), Inches(4.1), Inches(12), Inches(0.5),
             font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

    add_text(sl,
             "\"Empowering farmers with AI-driven price intelligence\"",
             Inches(0.5), Inches(5.0), Inches(12), Inches(0.6),
             font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_footer(sl, "AgriPredict AI  |  MCA Final Year Project  |  2026")
    return sl


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════

def create_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_title(prs)
    slide_problem(prs)
    slide_architecture(prs)
    slide_ml(prs)
    slide_tech(prs)
    slide_features(prs)
    slide_modules(prs)
    slide_results(prs)
    slide_database(prs)
    slide_multilingual(prs)
    slide_advisory(prs)
    slide_conclusion(prs)
    slide_thankyou(prs)

    save_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Project_Presentation.pptx")
    prs.save(save_path)
    print(f"✅  PPT saved: {save_path}")
    return save_path


if __name__ == "__main__":
    ppt_path = create_presentation()
    print(f"\n📄  File: {ppt_path}")
    print(f"📊  Slides: 13")
    print(f"\nTo convert to PDF, run:")
    print(f"  python generate_pdf.py")
